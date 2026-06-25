from io import BytesIO
import logging
import os
from pathlib import Path
import shutil
import subprocess
import tempfile
import time

from django.http import HttpResponse
from django.shortcuts import get_object_or_404, redirect, render
from django.urls import reverse
from django.views.decorators.http import require_POST
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from .forms import (
    AgencyNameFormSet,
    AgencyPageForm,
    BulletPageForm,
    CareerOpportunityPageForm,
    CaptionSlidePageForm,
    ClientDatePageForm,
    ClientSaysPageForm,
    FirstPageForm,
    GalleryPageForm,
    KudosBulletMomentPageForm,
    KudosMomentPageForm,
    NewJoinersPageForm,
    NextPageCardFormSet,
    NextPageForm,
    PromotionsPageForm,
    SpecialAppreciationPageForm,
)
from .models import (
    AgencyName,
    AgencyPage,
    BulletCard,
    BulletItem,
    BulletPage,
    CareerOpportunityJob,
    CareerOpportunityPage,
    CaptionSlideImage,
    CaptionSlidePage,
    ClientDateItem,
    ClientDatePage,
    ClientSaysPage,
    ClientSaysPerson,
    FirstPage,
    GalleryImage,
    GalleryPage,
    HiddenSlide,
    KudosBulletMomentCard,
    KudosBulletMomentPage,
    KudosMomentCard,
    KudosMomentPage,
    NewJoiner,
    NewJoinersPage,
    NextPage,
    NextPageCard,
    Promotion,
    PromotionsPage,
    SlideOrder,
    SpecialAppreciationPage,
)


logger = logging.getLogger(__name__)


PAGE_DELETE_CONFIG = {
    "first": (FirstPage, "dashboard"),
    "next": (NextPage, "next_page_dashboard"),
    "agency": (AgencyPage, "agency_page_dashboard"),
    "bullet": (BulletPage, "bullet_page_dashboard"),
    "client-date": (ClientDatePage, "client_date_page_dashboard"),
    "client-says": (ClientSaysPage, "client_says_page_dashboard"),
    "gallery": (GalleryPage, "gallery_page_dashboard"),
    "caption-slide": (CaptionSlidePage, "caption_slide_page_dashboard"),
    "new-joiners": (NewJoinersPage, "new_joiners_dashboard"),
    "promotions": (PromotionsPage, "promotions_dashboard"),
    "kudos-moments": (KudosMomentPage, "kudos_moments_dashboard"),
    "kudos-bullet-moments": (KudosBulletMomentPage, "kudos_bullet_moments_dashboard"),
    "special-appreciation": (SpecialAppreciationPage, "special_appreciation_dashboard"),
    "career-opportunities": (CareerOpportunityPage, "career_opportunities_dashboard"),
}


ITEM_DELETE_CONFIG = {
    "next-card": (NextPageCard, "next_page_dashboard"),
    "agency-name": (AgencyName, "agency_page_dashboard"),
    "bullet-card": (BulletCard, "bullet_page_dashboard"),
    "bullet-item": (BulletItem, "bullet_page_dashboard"),
    "client-date-item": (ClientDateItem, "client_date_page_dashboard"),
    "client-says-person": (ClientSaysPerson, "client_says_page_dashboard"),
    "gallery-image": (GalleryImage, "gallery_page_dashboard"),
    "caption-slide-image": (CaptionSlideImage, "caption_slide_page_dashboard"),
    "new-joiner": (NewJoiner, "new_joiners_dashboard"),
    "promotion": (Promotion, "promotions_dashboard"),
    "kudos-moment-card": (KudosMomentCard, "kudos_moments_dashboard"),
    "kudos-bullet-moment-card": (KudosBulletMomentCard, "kudos_bullet_moments_dashboard"),
    "career-opportunity-job": (CareerOpportunityJob, "career_opportunities_dashboard"),
}


SLIDE_PAGE_CONFIG = [
    ("first", FirstPage, "first_page_detail", "title"),
    ("next", NextPage, "next_page_view_detail", "title"),
    ("agency", AgencyPage, "agency_page_view_detail", "title"),
    ("bullet", BulletPage, "bullet_page_view_detail", "title"),
    ("client-date", ClientDatePage, "client_date_page_view_detail", "title"),
    ("client-says", ClientSaysPage, "client_says_page_view_detail", "heading"),
    ("gallery", GalleryPage, "gallery_page_view_detail", "caption"),
    ("caption-slide", CaptionSlidePage, "caption_slide_page_view_detail", "caption"),
    ("new-joiners", NewJoinersPage, "new_joiners_view_detail", "title"),
    ("promotions", PromotionsPage, "promotions_view_detail", "title"),
    ("kudos-moments", KudosMomentPage, "kudos_moments_view_detail", "title"),
    ("kudos-bullet-moments", KudosBulletMomentPage, "kudos_bullet_moments_view_detail", "title"),
    ("special-appreciation", SpecialAppreciationPage, "special_appreciation_view_detail", "title"),
    ("career-opportunities", CareerOpportunityPage, "career_opportunities_view_detail", "title"),
]


@require_POST
def delete_page_item(request, page_type, pk):
    model_redirect = PAGE_DELETE_CONFIG.get(page_type)
    if not model_redirect:
        return redirect("dashboard")

    model, redirect_name = model_redirect
    get_object_or_404(model, pk=pk).delete()
    next_url = request.POST.get("next")
    return redirect(next_url or redirect_name)


@require_POST
def delete_child_item(request, item_type, pk):
    model_redirect = ITEM_DELETE_CONFIG.get(item_type)
    if not model_redirect:
        return redirect("dashboard")

    model, redirect_name = model_redirect
    get_object_or_404(model, pk=pk).delete()
    next_url = request.POST.get("next")
    return redirect(next_url or redirect_name)


def all_slides_view(request):
    slides = get_visible_ordered_slides()
    return render(request, "decks/all_slides.html", {"slides": slides})


def get_visible_ordered_slides(order_keys=None):
    hidden_slides = set(HiddenSlide.objects.values_list("slide_type", "object_id"))
    saved_orders = {
        (order.slide_type, order.object_id): order.position
        for order in SlideOrder.objects.all()
    }
    slides = []
    default_position = 1

    for slide_type, model, detail_url_name, title_field in SLIDE_PAGE_CONFIG:
        for item in model.objects.all():
            if (slide_type, item.pk) in hidden_slides:
                continue

            slide_key = (slide_type, item.pk)
            slides.append(
                {
                    "type": slide_type,
                    "id": item.pk,
                    "item": item,
                    "title": getattr(item, title_field, str(item)),
                    "detail_url_name": detail_url_name,
                    "position": saved_orders.get(slide_key),
                    "default_position": default_position,
                }
            )
            default_position += 1

    if order_keys:
        requested_order = {key: index for index, key in enumerate(order_keys)}
        slides.sort(
            key=lambda slide: (
                requested_order.get(f"{slide['type']}:{slide['id']}", len(requested_order) + slide["default_position"])
            )
        )
    else:
        slides.sort(key=lambda slide: (slide["position"] is None, slide["position"] or slide["default_position"]))
    return slides


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color="0B1F35"):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text or "")
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor.from_string(color)
    return shape


def add_bullets(slide, left, top, width, height, items, font_size=16):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(item)
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = RGBColor.from_string("334155")
    return shape


def add_image_fit(slide, image_field, left, top, width, height):
    if not image_field:
        return

    try:
        slide.shapes.add_picture(image_field.path, left, top, width=width, height=height)
    except (OSError, ValueError):
        return


def add_title(slide, title, subtitle=""):
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.5), title, 28, True)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.9), Inches(11.8), Inches(0.35), subtitle, 13, False, "64748B")


def render_slide_to_ppt(prs, slide_data):
    item = slide_data["item"]
    slide_type = slide_data["type"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string("F8FAFC")

    if slide_type == "first":
        add_image_fit(slide, item.image, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        add_textbox(slide, Inches(0.75), Inches(4.75), Inches(8.4), Inches(0.9), item.title, 38, True, "FFFFFF")
        add_textbox(slide, Inches(0.78), Inches(5.7), Inches(8.5), Inches(0.5), item.description, 16, False, "FFFFFF")
        add_textbox(slide, Inches(10.1), Inches(5.65), Inches(2.3), Inches(0.5), f"{item.get_month_display()} {item.year}", 17, True, "FFFFFF")
        return

    subtitle = getattr(item, "caption", "") or getattr(item, "description", "") or getattr(item, "client_name", "")
    add_title(slide, slide_data["title"], subtitle)

    if slide_type == "next":
        add_textbox(slide, Inches(0.7), Inches(1.55), Inches(2.4), Inches(0.75), item.score_value, 32, True, "082A74")
        add_textbox(slide, Inches(0.75), Inches(2.25), Inches(2.4), Inches(0.35), item.score_title, 13, True, "64748B")
        for idx, card in enumerate(item.cards.all()):
            x = Inches(3.4 + (idx % 3) * 3.1)
            y = Inches(1.55 + (idx // 3) * 1.3)
            add_textbox(slide, x, y, Inches(2.6), Inches(0.45), card.value, 25, True, "082A74")
            add_textbox(slide, x, y + Inches(0.48), Inches(2.6), Inches(0.35), card.label, 12, False, "64748B")
    elif slide_type == "agency":
        add_textbox(slide, Inches(0.75), Inches(1.45), Inches(5.2), Inches(1.0), item.content, 18)
        add_bullets(slide, Inches(6.3), Inches(1.45), Inches(5.7), Inches(4.8), [agency.name for agency in item.agencies.all()], 18)
    elif slide_type == "bullet":
        y = Inches(1.4)
        for card in item.cards.all():
            add_textbox(slide, Inches(0.75), y, Inches(4.2), Inches(0.35), card.title, 17, True, "082A74")
            add_bullets(slide, Inches(1.0), y + Inches(0.4), Inches(5.2), Inches(1.0), [bullet.text for bullet in card.items.all()], 13)
            y += Inches(1.3)
    elif slide_type == "client-date":
        add_textbox(slide, Inches(0.75), Inches(1.45), Inches(5.3), Inches(1.1), item.content, 17)
        add_bullets(slide, Inches(6.35), Inches(1.45), Inches(5.7), Inches(4.8), [f"{client.name} - {client.went_live_on:%m/%d/%y}" for client in item.clients.all()], 16)
    elif slide_type == "client-says":
        add_textbox(slide, Inches(0.85), Inches(1.55), Inches(6.2), Inches(2.2), f'"{item.client_description}"', 21, False, "0B1F35")
        add_textbox(slide, Inches(0.85), Inches(3.85), Inches(5.5), Inches(0.4), item.client_name, 16, True, "082A74")
        for idx, person in enumerate(item.people.all()[:6]):
            x = Inches(7.2 + (idx % 3) * 1.65)
            y = Inches(1.45 + (idx // 3) * 2.05)
            add_image_fit(slide, person.image, x, y, Inches(1.15), Inches(1.15))
            add_textbox(slide, x, y + Inches(1.2), Inches(1.5), Inches(0.3), person.name, 9, True)
            add_textbox(slide, x, y + Inches(1.5), Inches(1.5), Inches(0.35), person.designation, 8, False, "64748B")
    elif slide_type == "gallery":
        for idx, image in enumerate(item.images.all()[:5]):
            x = Inches(0.75 + (idx % 3) * 4.0)
            y = Inches(1.45 + (idx // 3) * 2.45)
            add_image_fit(slide, image.image, x, y, Inches(3.45), Inches(2.0))
    elif slide_type == "caption-slide":
        for idx, image in enumerate(item.images.all()[:3]):
            x = Inches(0.75 + idx * 4.1)
            add_image_fit(slide, image.image, x, Inches(1.55), Inches(3.55), Inches(4.25))
    elif slide_type in {"new-joiners", "promotions"}:
        people = item.joiners.all() if slide_type == "new-joiners" else item.promotions.all()
        for idx, person in enumerate(people[:10]):
            x = Inches(0.75 + (idx % 5) * 2.45)
            y = Inches(1.35 + (idx // 5) * 2.65)
            add_image_fit(slide, person.image, x, y, Inches(1.35), Inches(1.35))
            add_textbox(slide, x, y + Inches(1.45), Inches(2.1), Inches(0.3), person.name, 10, True)
            add_textbox(slide, x, y + Inches(1.75), Inches(2.1), Inches(0.35), person.designation, 8, False, "64748B")
    elif slide_type in {"kudos-moments", "kudos-bullet-moments"}:
        if item.image:
            add_image_fit(slide, item.image, Inches(0.8), Inches(1.45), Inches(3.0), Inches(3.4))
        if item.name:
            add_textbox(slide, Inches(0.9), Inches(5.0), Inches(3.0), Inches(0.35), item.name, 14, True)
            add_textbox(slide, Inches(0.9), Inches(5.35), Inches(3.0), Inches(0.35), item.designation, 10, False, "64748B")
        for idx, card in enumerate(item.cards.all()[:10]):
            x = Inches(4.3 + (idx % 3) * 2.8)
            y = Inches(1.45 + (idx // 3) * 1.15)
            headline = getattr(card, "value", "") or getattr(card, "title", "")
            detail = getattr(card, "label", "") or getattr(card, "content", "")
            add_textbox(slide, x, y, Inches(2.45), Inches(0.35), headline, 18, True, "082A74")
            add_textbox(slide, x, y + Inches(0.38), Inches(2.45), Inches(0.45), detail, 9, False, "64748B")
    elif slide_type == "special-appreciation":
        add_image_fit(slide, item.image, Inches(0.75), Inches(1.45), Inches(3.2), Inches(3.8))
        add_textbox(slide, Inches(4.35), Inches(1.45), Inches(2.8), Inches(0.35), item.name, 17, True, "082A74")
        add_textbox(slide, Inches(4.35), Inches(1.85), Inches(2.8), Inches(0.35), item.designation, 11, False, "64748B")
        add_textbox(slide, Inches(4.35), Inches(2.35), Inches(7.1), Inches(2.2), item.content, 17)
        if item.author:
            add_textbox(slide, Inches(4.35), Inches(4.8), Inches(4.0), Inches(0.35), f"- {item.author}", 13, True, "082A74")
    elif slide_type == "career-opportunities":
        add_textbox(slide, Inches(0.75), Inches(1.45), Inches(5.0), Inches(1.2), item.content, 17)
        add_bullets(slide, Inches(6.15), Inches(1.45), Inches(5.9), Inches(4.8), [job.title for job in item.jobs.all()], 18)


def edge_executable():
    candidates = [
        shutil.which("msedge"),
        shutil.which("chrome"),
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return candidate
    return None


def capture_slide_screenshot_with_playwright(url, output_path):
    os.environ.setdefault("PLAYWRIGHT_BROWSERS_PATH", "0")
    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        logger.warning("Playwright is not installed: %s", exc)
        return False, f"Playwright is not installed: {exc}"

    try:
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch(
                args=[
                    "--disable-dev-shm-usage",
                    "--disable-gpu",
                    "--no-sandbox",
                ]
            )
            page = browser.new_page(viewport={"width": 1200, "height": 675}, device_scale_factor=1)
            response = page.goto(url, wait_until="domcontentloaded", timeout=10000)
            if response and not response.ok:
                raise RuntimeError(f"Slide URL returned HTTP {response.status}: {url}")
            try:
                page.wait_for_load_state("networkidle", timeout=5000)
            except Exception:
                logger.info("Continuing screenshot before network idle for %s.", url)
            page.wait_for_function(
                "() => Array.from(document.images).every((image) => image.complete)",
                timeout=5000,
            )
            page.screenshot(path=str(output_path), full_page=False)
            browser.close()
    except Exception as exc:
        logger.warning("Playwright screenshot failed for %s: %s", url, exc)
        return False, f"Playwright screenshot failed: {exc}"

    if output_path.exists() and output_path.stat().st_size > 0:
        return True, ""
    return False, f"Playwright did not create a screenshot file: {output_path}"


def capture_slide_screenshot_with_edge(url, output_path, user_data_dir):
    browser = edge_executable()
    if not browser:
        return False, "No Edge or Chrome executable was found."

    command = [
        browser,
        "--headless=new",
        "--disable-gpu",
        "--hide-scrollbars",
        "--force-device-scale-factor=1",
        f"--user-data-dir={user_data_dir}",
        "--window-size=1200,675",
        f"--screenshot={output_path}",
        url,
    ]
    try:
        result = subprocess.run(command, capture_output=True, timeout=10)
    except (OSError, subprocess.SubprocessError) as exc:
        logger.warning("Browser screenshot failed to start for %s: %s", url, exc)
        return False, f"Browser screenshot failed to start: {exc}"

    if result.returncode != 0:
        stderr = result.stderr.decode(errors="ignore")
        logger.warning("Browser screenshot failed for %s: %s", url, stderr)
        return False, f"Browser screenshot failed: {stderr}"

    if Path(output_path).exists() and Path(output_path).stat().st_size > 0:
        return True, ""
    return False, f"Browser did not create a screenshot file: {output_path}"


def capture_slide_screenshot(url, output_path, user_data_dir):
    success, error = capture_slide_screenshot_with_playwright(url, output_path)
    if success:
        return True, ""

    if os.name == "nt":
        edge_success, edge_error = capture_slide_screenshot_with_edge(url, output_path, user_data_dir)
        if edge_success:
            return True, ""
        return False, f"{error}; {edge_error}"

    return False, error


def render_slide_screenshot_to_ppt(prs, screenshot_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(screenshot_path, 0, 0, width=prs.slide_width, height=prs.slide_height)


def build_slide_capture_url(request, detail_url_name, object_id):
    path = reverse(detail_url_name, args=[object_id])
    port = os.environ.get("PORT")
    if port:
        return f"http://127.0.0.1:{port}{path}"
    return request.build_absolute_uri(path)


def render_export_error_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string("F8FAFC")
    add_textbox(slide, Inches(0.75), Inches(2.8), Inches(11.8), Inches(0.5), title, 22, True)
    add_textbox(
        slide,
        Inches(0.75),
        Inches(3.35),
        Inches(11.8),
        Inches(0.45),
        "This slide could not be rendered during export.",
        14,
        False,
        "64748B",
    )


def remove_temp_tree_later(path):
    for _ in range(5):
        try:
            shutil.rmtree(path)
            return
        except (OSError, PermissionError):
            time.sleep(0.25)
    shutil.rmtree(path, ignore_errors=True)


def powerpoint_response(prs, filename="ordered-slides.pptx"):
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    response = HttpResponse(
        output.getvalue(),
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    response["Content-Disposition"] = f'attachment; filename="{filename}"'
    return response


def export_slides_powerpoint(request):
    temp_dir = None
    browser_profile_dir = None
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        order_keys = request.POST.getlist("slide_key")
        slides = get_visible_ordered_slides(order_keys=order_keys)
        use_screenshots = request.get_host() != "testserver"

        temp_dir = tempfile.mkdtemp()
        browser_profile_dir = tempfile.mkdtemp()
        screenshots_available = use_screenshots
        for index, slide_data in enumerate(slides, start=1):
            slide_url = build_slide_capture_url(request, slide_data["detail_url_name"], slide_data["id"])
            screenshot_path = Path(temp_dir) / f"slide-{index}.png"
            try:
                if screenshots_available:
                    captured, capture_error = capture_slide_screenshot(slide_url, screenshot_path, browser_profile_dir)
                    if captured:
                        render_slide_screenshot_to_ppt(prs, str(screenshot_path))
                    else:
                        screenshots_available = False
                        logger.error(
                            "Designed PowerPoint export failed for %s. Falling back to generated slides for this export. Reason: %s",
                            slide_url,
                            capture_error,
                        )
                        render_slide_to_ppt(prs, slide_data)
                else:
                    render_slide_to_ppt(prs, slide_data)
            except Exception:
                logger.exception("Designed PowerPoint export failed for slide %s.", slide_data)
                raise

        return powerpoint_response(prs)
    except Exception as exc:
        logger.exception("PowerPoint export failed before response could be created.")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        render_export_error_slide(prs, "PowerPoint export failed")
        add_textbox(
            prs.slides[-1],
            Inches(0.75),
            Inches(3.85),
            Inches(11.8),
            Inches(0.45),
            str(exc),
            10,
            False,
            "B42318",
        )
        return powerpoint_response(prs, filename="export-error.pptx")
    finally:
        if temp_dir:
            remove_temp_tree_later(temp_dir)
        if browser_profile_dir:
            remove_temp_tree_later(browser_profile_dir)


@require_POST
def remove_slide_from_view(request):
    slide_key = request.POST.get("slide_key", "")
    valid_keys = {slide_type for slide_type, _, _, _ in SLIDE_PAGE_CONFIG}
    if ":" in slide_key and slide_key.split(":", 1)[0] in valid_keys:
        slide_type, object_id = slide_key.split(":", 1)
        if object_id.isdigit():
            HiddenSlide.objects.get_or_create(slide_type=slide_type, object_id=int(object_id))

    return redirect("all_slides")


@require_POST
def save_slide_order(request):
    valid_keys = {slide_type for slide_type, _, _, _ in SLIDE_PAGE_CONFIG}
    slide_keys = request.POST.getlist("slide_key")

    for position, slide_key in enumerate(slide_keys, start=1):
        if ":" not in slide_key:
            continue

        slide_type, object_id = slide_key.split(":", 1)
        if slide_type not in valid_keys or not object_id.isdigit():
            continue

        SlideOrder.objects.update_or_create(
            slide_type=slide_type,
            object_id=int(object_id),
            defaults={"position": position},
        )

    return redirect("all_slides")


@require_POST
def reset_removed_slides(request):
    HiddenSlide.objects.all().delete()
    return redirect("all_slides")


def first_page_list(request, pk=None):
    first_pages = FirstPage.objects.all()
    selected_page = get_object_or_404(first_pages, pk=pk) if pk else first_pages.first()
    return render(
        request,
        "decks/first_page_list.html",
        {"first_pages": first_pages, "selected_page": selected_page},
    )


def next_page_view(request, pk=None):
    queryset = NextPage.objects.prefetch_related("cards")
    next_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    return render(request, "decks/next_page_view.html", {"next_page": next_page})


def agency_page_view(request, pk=None):
    queryset = AgencyPage.objects.prefetch_related("agencies")
    agency_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    agencies = agency_page.agencies.all() if agency_page else []
    agency_image = "decks/img/image-7.png"
    if agency_page and "churn" in agency_page.title.lower():
        agency_image = "decks/img/image-8.png"
    return render(
        request,
        "decks/agency_page_view.html",
        {
            "agency_page": agency_page,
            "agencies": agencies,
            "agency_image": agency_image,
        },
    )


def bullet_page_view(request, pk=None):
    queryset = BulletPage.objects.prefetch_related("cards__items")
    bullet_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    return render(request, "decks/bullet_page_view.html", {"bullet_page": bullet_page})


def client_date_page_view(request, pk=None):
    queryset = ClientDatePage.objects.prefetch_related("clients")
    client_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    return render(request, "decks/client_date_page_view.html", {"client_page": client_page})


def client_says_page_view(request, pk=None):
    queryset = ClientSaysPage.objects.prefetch_related("people")
    says_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    people = list(says_page.people.all()) if says_page else []
    team_leader = next((person for person in people if person.is_team_leader), None)
    regular_people = [person for person in people if person != team_leader]
    use_people_grid = (
        len(regular_people) > 5
        if team_leader is None
        else len(regular_people) >= 5
    )
    return render(
        request,
        "decks/client_says_page_view.html",
        {
            "says_page": says_page,
            "team_leader": team_leader,
            "regular_people": regular_people,
            "people_count": len(regular_people),
            "total_people_count": len(people),
            "use_people_grid": use_people_grid,
        },
    )


def gallery_page_view(request, pk=None):
    queryset = GalleryPage.objects.prefetch_related("images")
    gallery_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    return render(request, "decks/gallery_page_view.html", {"gallery_page": gallery_page})


def caption_slide_page_view(request, pk=None):
    queryset = CaptionSlidePage.objects.prefetch_related("images")
    caption_slide_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    return render(
        request,
        "decks/caption_slide_page_view.html",
        {"caption_slide_page": caption_slide_page},
    )


def new_joiners_view(request, pk=None):
    queryset = NewJoinersPage.objects.prefetch_related("joiners")
    new_joiners_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    return render(
        request,
        "decks/new_joiners_view.html",
        {"new_joiners_page": new_joiners_page},
    )


def promotions_view(request, pk=None):
    queryset = PromotionsPage.objects.prefetch_related("promotions")
    promotions_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    return render(
        request,
        "decks/promotions_view.html",
        {"promotions_page": promotions_page},
    )


def kudos_moments_view(request, pk=None):
    queryset = KudosMomentPage.objects.prefetch_related("cards")
    kudos_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    return render(
        request,
        "decks/kudos_moments_view.html",
        {"kudos_page": kudos_page},
    )


def kudos_bullet_moments_view(request, pk=None):
    queryset = KudosBulletMomentPage.objects.prefetch_related("cards")
    kudos_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    return render(
        request,
        "decks/kudos_bullet_moments_view.html",
        {"kudos_page": kudos_page},
    )


def special_appreciation_view(request, pk=None):
    queryset = SpecialAppreciationPage.objects.all()
    special_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    return render(
        request,
        "decks/special_appreciation_view.html",
        {"special_page": special_page},
    )


def career_opportunities_view(request, pk=None):
    queryset = CareerOpportunityPage.objects.prefetch_related("jobs")
    career_page = get_object_or_404(queryset, pk=pk) if pk else queryset.first()
    jobs = list(career_page.jobs.all()) if career_page else []
    show_career_image_outside = len(jobs) > 3 or any(
        len(job.title) > 28 or len(job.content) > 55 for job in jobs
    )
    return render(
        request,
        "decks/career_opportunities_view.html",
        {
            "career_page": career_page,
            "career_jobs": jobs,
            "show_career_image_outside": show_career_image_outside,
        },
    )


def dashboard(request):
    if request.method == "POST":
        form = FirstPageForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            return redirect("first_page_list")
    else:
        form = FirstPageForm()

    first_pages = FirstPage.objects.all()
    return render(
        request,
        "decks/dashboard.html",
        {"form": form, "first_pages": first_pages, "active_menu": "first"},
    )


def next_page_dashboard(request):
    page = None

    if request.method == "POST":
        form = NextPageForm(request.POST)
        formset = NextPageCardFormSet(request.POST)
        if form.is_valid():
            page = form.save()
            formset = NextPageCardFormSet(request.POST, instance=page)
            if formset.is_valid():
                formset.save()
                return redirect("next_page_dashboard")
    else:
        form = NextPageForm()
        formset = NextPageCardFormSet()

    next_pages = NextPage.objects.prefetch_related("cards").all()
    return render(
        request,
        "decks/next_page_dashboard.html",
        {
            "form": form,
            "formset": formset,
            "next_pages": next_pages,
            "active_menu": "next",
        },
    )


def agency_page_dashboard(request):
    page = None

    if request.method == "POST":
        form = AgencyPageForm(request.POST)
        formset = AgencyNameFormSet(request.POST)
        if form.is_valid():
            page = form.save()
            formset = AgencyNameFormSet(request.POST, instance=page)
            if formset.is_valid():
                formset.save()
                return redirect("agency_page_dashboard")
    else:
        form = AgencyPageForm()
        formset = AgencyNameFormSet()

    agency_pages = AgencyPage.objects.prefetch_related("agencies").all()
    return render(
        request,
        "decks/agency_page_dashboard.html",
        {
            "form": form,
            "formset": formset,
            "agency_pages": agency_pages,
            "active_menu": "agency",
        },
    )


def bullet_page_dashboard(request):
    page = None

    if request.method == "POST":
        form = BulletPageForm(request.POST)
        if form.is_valid():
            page = form.save()

            card_titles = request.POST.getlist("card_title")
            bullet_groups = request.POST.getlist("card_bullets")
            for index, card_title in enumerate(card_titles, start=1):
                title = card_title.strip()
                if not title:
                    continue

                card = BulletCard.objects.create(page=page, title=title, order=index)
                bullets = bullet_groups[index - 1].splitlines() if index - 1 < len(bullet_groups) else []
                for bullet_index, bullet in enumerate(bullets, start=1):
                    text = bullet.strip()
                    if text:
                        BulletItem.objects.create(card=card, text=text, order=bullet_index)

            return redirect("bullet_page_dashboard")
    else:
        form = BulletPageForm()

    bullet_pages = BulletPage.objects.prefetch_related("cards__items").all()
    return render(
        request,
        "decks/bullet_page_dashboard.html",
        {
            "form": form,
            "page": page,
            "bullet_pages": bullet_pages,
            "active_menu": "bullet",
        },
    )


def client_date_page_dashboard(request):
    page = None

    if request.method == "POST":
        form = ClientDatePageForm(request.POST)
        if form.is_valid():
            page = form.save()

            names = request.POST.getlist("client_name")
            dates = request.POST.getlist("went_live_on")
            orders = request.POST.getlist("client_order")
            for index, name in enumerate(names, start=1):
                client_name = name.strip()
                went_live_on = dates[index - 1] if index - 1 < len(dates) else ""
                order = orders[index - 1] if index - 1 < len(orders) else index
                if client_name and went_live_on:
                    ClientDateItem.objects.create(
                        page=page,
                        name=client_name,
                        went_live_on=went_live_on,
                        order=order or index,
                    )

            return redirect("client_date_page_dashboard")
    else:
        form = ClientDatePageForm()

    client_pages = ClientDatePage.objects.prefetch_related("clients").all()
    return render(
        request,
        "decks/client_date_page_dashboard.html",
        {
            "form": form,
            "page": page,
            "client_pages": client_pages,
            "active_menu": "client_date",
        },
    )


def client_says_page_dashboard(request):
    page = None

    if request.method == "POST":
        form = ClientSaysPageForm(request.POST)
        if form.is_valid():
            page = form.save()
            seen_people = []

            ids = request.POST.getlist("person_id")
            names = request.POST.getlist("person_name")
            designations = request.POST.getlist("person_designation")
            orders = request.POST.getlist("person_order")

            for index, name in enumerate(names, start=1):
                person_name = name.strip()
                designation = designations[index - 1].strip() if index - 1 < len(designations) else ""
                order = orders[index - 1] if index - 1 < len(orders) else index
                person_id = ids[index - 1] if index - 1 < len(ids) else ""
                image = request.FILES.get(f"person_image_{index}")

                if not person_name or not designation:
                    continue

                person = ClientSaysPerson(page=page)
                person.name = person_name
                person.designation = designation
                person.order = order or index
                if image:
                    person.image = image
                if person.image:
                    person.save()
                    seen_people.append(person.pk)

            return redirect("client_says_page_dashboard")
    else:
        form = ClientSaysPageForm()

    says_pages = ClientSaysPage.objects.prefetch_related("people").all()
    return render(
        request,
        "decks/client_says_page_dashboard.html",
        {
            "form": form,
            "page": page,
            "says_pages": says_pages,
            "active_menu": "client_says",
        },
    )


def gallery_page_dashboard(request):
    page = None

    if request.method == "POST":
        form = GalleryPageForm(request.POST)
        if form.is_valid():
            page = form.save()
            seen_images = []

            ids = request.POST.getlist("gallery_image_id")
            orders = request.POST.getlist("gallery_image_order")

            for index in range(5):
                image_id = ids[index] if index < len(ids) else ""
                order = orders[index] if index < len(orders) else index + 1
                upload = request.FILES.get(f"gallery_image_{index + 1}")

                gallery_image = GalleryImage(page=page)
                gallery_image.order = order or index + 1
                if upload:
                    gallery_image.image = upload
                if gallery_image.image:
                    gallery_image.save()
                    seen_images.append(gallery_image.pk)

            return redirect("gallery_page_dashboard")
    else:
        form = GalleryPageForm()

    gallery_pages = GalleryPage.objects.prefetch_related("images").all()
    gallery_image_rows = [None] * 5
    return render(
        request,
        "decks/gallery_page_dashboard.html",
        {
            "form": form,
            "page": page,
            "gallery_image_rows": gallery_image_rows,
            "gallery_pages": gallery_pages,
            "active_menu": "gallery",
        },
    )


def caption_slide_page_dashboard(request):
    page = None

    if request.method == "POST":
        form = CaptionSlidePageForm(request.POST)
        if form.is_valid():
            page = form.save()
            seen_images = []

            ids = request.POST.getlist("caption_slide_image_id")
            orders = request.POST.getlist("caption_slide_image_order")

            for index in range(3):
                image_id = ids[index] if index < len(ids) else ""
                order = orders[index] if index < len(orders) else index + 1
                upload = request.FILES.get(f"caption_slide_image_{index + 1}")

                slide_image = CaptionSlideImage(page=page)
                slide_image.order = order or index + 1
                if upload:
                    slide_image.image = upload
                if slide_image.image:
                    slide_image.save()
                    seen_images.append(slide_image.pk)

            return redirect("caption_slide_page_dashboard")
    else:
        form = CaptionSlidePageForm()

    caption_slide_pages = CaptionSlidePage.objects.prefetch_related("images").all()
    caption_slide_image_rows = [None] * 3
    return render(
        request,
        "decks/caption_slide_page_dashboard.html",
        {
            "form": form,
            "page": page,
            "caption_slide_image_rows": caption_slide_image_rows,
            "caption_slide_pages": caption_slide_pages,
            "active_menu": "caption_slide",
        },
    )


def new_joiners_dashboard(request):
    page = None

    if request.method == "POST":
        form = NewJoinersPageForm(request.POST)
        if form.is_valid():
            page = form.save()
            seen_joiners = []

            ids = request.POST.getlist("new_joiner_id")
            names = request.POST.getlist("new_joiner_name")
            designations = request.POST.getlist("new_joiner_designation")
            orders = request.POST.getlist("new_joiner_order")

            for index in range(10):
                joiner_id = ids[index] if index < len(ids) else ""
                name = names[index].strip() if index < len(names) else ""
                designation = designations[index].strip() if index < len(designations) else ""
                order = orders[index] if index < len(orders) else index + 1
                upload = request.FILES.get(f"new_joiner_image_{index + 1}")

                joiner = NewJoiner(page=page)
                joiner.name = name
                joiner.designation = designation
                joiner.order = order or index + 1
                if upload:
                    joiner.image = upload
                if joiner.name and joiner.designation and joiner.image:
                    joiner.save()
                    seen_joiners.append(joiner.pk)

            return redirect("new_joiners_dashboard")
    else:
        form = NewJoinersPageForm()

    new_joiners_pages = NewJoinersPage.objects.prefetch_related("joiners").all()
    new_joiner_rows = [None] * 10
    return render(
        request,
        "decks/new_joiners_dashboard.html",
        {
            "form": form,
            "page": page,
            "new_joiner_rows": new_joiner_rows,
            "new_joiners_pages": new_joiners_pages,
            "active_menu": "new_joiners",
        },
    )


def promotions_dashboard(request):
    page = None

    if request.method == "POST":
        form = PromotionsPageForm(request.POST)
        if form.is_valid():
            page = form.save()
            seen_promotions = []

            ids = request.POST.getlist("promotion_id")
            names = request.POST.getlist("promotion_name")
            designations = request.POST.getlist("promotion_designation")
            orders = request.POST.getlist("promotion_order")

            for index in range(10):
                promotion_id = ids[index] if index < len(ids) else ""
                name = names[index].strip() if index < len(names) else ""
                designation = designations[index].strip() if index < len(designations) else ""
                order = orders[index] if index < len(orders) else index + 1
                upload = request.FILES.get(f"promotion_image_{index + 1}")

                promotion = Promotion(page=page)
                promotion.name = name
                promotion.designation = designation
                promotion.order = order or index + 1
                if upload:
                    promotion.image = upload
                if promotion.name and promotion.designation and promotion.image:
                    promotion.save()
                    seen_promotions.append(promotion.pk)

            return redirect("promotions_dashboard")
    else:
        form = PromotionsPageForm()

    promotions_pages = PromotionsPage.objects.prefetch_related("promotions").all()
    promotion_rows = [None] * 10
    return render(
        request,
        "decks/promotions_dashboard.html",
        {
            "form": form,
            "page": page,
            "promotion_rows": promotion_rows,
            "promotions_pages": promotions_pages,
            "active_menu": "promotions",
        },
    )


def kudos_moments_dashboard(request):
    page = None

    if request.method == "POST":
        form = KudosMomentPageForm(request.POST, request.FILES)
        if form.is_valid():
            page = form.save()
            seen_cards = []

            ids = request.POST.getlist("kudos_card_id")
            values = request.POST.getlist("kudos_card_value")
            labels = request.POST.getlist("kudos_card_label")
            orders = request.POST.getlist("kudos_card_order")

            for index in range(10):
                card_id = ids[index] if index < len(ids) else ""
                value = values[index].strip() if index < len(values) else ""
                label = labels[index].strip() if index < len(labels) else ""
                order = orders[index] if index < len(orders) else index + 1

                card = KudosMomentCard(page=page)
                card.value = value
                card.label = label
                card.order = order or index + 1
                if card.value and card.label:
                    card.save()
                    seen_cards.append(card.pk)

            return redirect("kudos_moments_dashboard")
    else:
        form = KudosMomentPageForm()

    kudos_pages = KudosMomentPage.objects.prefetch_related("cards").all()
    kudos_card_rows = [None] * 10
    return render(
        request,
        "decks/kudos_moments_dashboard.html",
        {
            "form": form,
            "page": page,
            "kudos_card_rows": kudos_card_rows,
            "kudos_pages": kudos_pages,
            "active_menu": "kudos_moments",
        },
    )


def special_appreciation_dashboard(request):
    page = None

    if request.method == "POST":
        form = SpecialAppreciationPageForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            return redirect("special_appreciation_dashboard")
    else:
        form = SpecialAppreciationPageForm()

    special_pages = SpecialAppreciationPage.objects.all()
    return render(
        request,
        "decks/special_appreciation_dashboard.html",
        {
            "form": form,
            "page": page,
            "special_pages": special_pages,
            "active_menu": "special_appreciation",
        },
    )


def kudos_bullet_moments_dashboard(request):
    page = None

    if request.method == "POST":
        form = KudosBulletMomentPageForm(request.POST, request.FILES)
        if form.is_valid():
            page = form.save()
            seen_cards = []

            ids = request.POST.getlist("kudos_bullet_card_id")
            card_types = request.POST.getlist("kudos_bullet_card_type")
            values = request.POST.getlist("kudos_bullet_card_value")
            labels = request.POST.getlist("kudos_bullet_card_label")
            titles = request.POST.getlist("kudos_bullet_card_title")
            contents = request.POST.getlist("kudos_bullet_card_content")
            orders = request.POST.getlist("kudos_bullet_card_order")

            for index in range(10):
                card_id = ids[index] if index < len(ids) else ""
                card_type = card_types[index] if index < len(card_types) else KudosBulletMomentCard.METRIC
                value = values[index].strip() if index < len(values) else ""
                label = labels[index].strip() if index < len(labels) else ""
                title = titles[index].strip() if index < len(titles) else ""
                content = contents[index].strip() if index < len(contents) else ""
                order = orders[index] if index < len(orders) else index + 1

                card = KudosBulletMomentCard(page=page)
                card.card_type = card_type
                card.value = value
                card.label = label
                card.title = title
                card.content = content
                card.order = order or index + 1

                if card.value or card.label or card.title or card.content:
                    card.save()
                    seen_cards.append(card.pk)

            return redirect("kudos_bullet_moments_dashboard")
    else:
        form = KudosBulletMomentPageForm()

    kudos_pages = KudosBulletMomentPage.objects.prefetch_related("cards").all()
    kudos_card_rows = [None] * 10
    card_type_choices = KudosBulletMomentCard.CARD_TYPE_CHOICES
    return render(
        request,
        "decks/kudos_bullet_moments_dashboard.html",
        {
            "form": form,
            "page": page,
            "kudos_card_rows": kudos_card_rows,
            "kudos_pages": kudos_pages,
            "card_type_choices": card_type_choices,
            "active_menu": "kudos_bullet_moments",
        },
    )


def career_opportunities_dashboard(request):
    page = None

    if request.method == "POST":
        form = CareerOpportunityPageForm(request.POST)
        if form.is_valid():
            page = form.save()
            seen_jobs = []

            ids = request.POST.getlist("career_job_id")
            titles = request.POST.getlist("career_job_title")
            contents = request.POST.getlist("career_job_content")
            orders = request.POST.getlist("career_job_order")

            for index in range(10):
                job_id = ids[index] if index < len(ids) else ""
                title = titles[index].strip() if index < len(titles) else ""
                content = contents[index].strip() if index < len(contents) else ""
                order = orders[index] if index < len(orders) else index + 1

                job = CareerOpportunityJob(page=page)
                job.title = title
                job.content = content
                job.order = order or index + 1
                if job.title:
                    job.save()
                    seen_jobs.append(job.pk)

            return redirect("career_opportunities_dashboard")
    else:
        form = CareerOpportunityPageForm()

    career_pages = CareerOpportunityPage.objects.prefetch_related("jobs").all()
    career_job_rows = [None] * 10
    return render(
        request,
        "decks/career_opportunities_dashboard.html",
        {
            "form": form,
            "page": page,
            "career_job_rows": career_job_rows,
            "career_pages": career_pages,
            "active_menu": "career_opportunities",
        },
    )
